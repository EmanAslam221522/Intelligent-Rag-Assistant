import os
import tempfile
from typing import Optional
import logging

# File processing libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import pandas as pd
except ImportError:
    pd = None

class FileProcessor:
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.doc': self._extract_docx,
            '.txt': self._extract_text,
            '.pptx': self._extract_pptx,
            '.csv': self._extract_csv,
            '.md': self._extract_text
        }
    
    def extract_text(self, file_path: str, filename: str) -> str:
        """Extract text from various file formats"""
        try:
            file_ext = os.path.splitext(filename.lower())[1]
            
            if file_ext not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            extractor = self.supported_formats[file_ext]
            return extractor(file_path)
            
        except Exception as e:
            logging.error(f"Error extracting text from {filename}: {str(e)}")
            raise Exception(f"Failed to extract text from {filename}: {str(e)}")
    
    def _extract_text(self, file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            raise Exception("Unable to decode text file")
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files"""
        if PyPDF2 is None:
            raise Exception("PyPDF2 library not installed. Install with: pip install PyPDF2")
        
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
        
        return text.strip()
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files"""
        if Document is None:
            raise Exception("python-docx library not installed. Install with: pip install python-docx")
        
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\n"
        except Exception as e:
            raise Exception(f"Error reading DOCX: {str(e)}")
        
        return text.strip()
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX files"""
        if pptx is None:
            raise Exception("python-pptx library not installed. Install with: pip install python-pptx")
        
        try:
            prs = pptx.Presentation(file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"Error reading PPTX: {str(e)}")
        
        return text.strip()
    
    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files"""
        if pd is None:
            raise Exception("pandas library not installed. Install with: pip install pandas")
        
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            raise Exception(f"Error reading CSV: {str(e)}")
    
    def get_file_info(self, file_path: str) -> dict:
        """Get basic file information"""
        try:
            stat = os.stat(file_path)
            return {
                "size": stat.st_size,
                "modified": stat.st_mtime,
                "created": stat.st_ctime
            }
        except Exception as e:
            return {"error": str(e)}
    
    def validate_file(self, filename: str, file_size: int) -> tuple[bool, str]:
        """Validate file before processing"""
        # Check file extension
        file_ext = os.path.splitext(filename.lower())[1]
        if file_ext not in self.supported_formats:
            return False, f"Unsupported file format: {file_ext}"
        
        # Check file size (limit to 10MB)
        max_size = 10 * 1024 * 1024  # 10MB
        if file_size > max_size:
            return False, f"File too large. Maximum size is {max_size // (1024*1024)}MB"
        
        return True, "Valid file"


